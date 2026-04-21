from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models import DeckSpec, SlideSpec
from ..utils.files import ensure_parent_dir
from .theme_registry import THEMES


class PPTXRenderer:
    def __init__(self, style_direction: str, total_slides: int = 0) -> None:
        self.theme = THEMES[style_direction]
        self.total_slides = total_slides
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _add_background(self, slide) -> None:
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            self.prs.slide_width,
            self.prs.slide_height,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme["background"]
        background.line.fill.background()
        sp_tree = slide.shapes._spTree
        element = background._element
        sp_tree.remove(element)
        sp_tree.insert(2, element)

    def _add_orb(self, slide, left: float, top: float, size: float, color: RGBColor, transparency: float) -> None:
        orb = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(size),
            Inches(size),
        )
        orb.fill.solid()
        orb.fill.fore_color.rgb = color
        orb.fill.transparency = transparency
        orb.line.fill.background()

    def _add_text(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int,
        color: RGBColor,
        font_name: str | None = None,
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        line_spacing: float = 1.12,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Pt(2)
        frame.margin_right = Pt(2)
        frame.margin_top = Pt(1)
        frame.margin_bottom = Pt(1)
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = font_name or self.theme["font"]
        return box

    def _fit_font(self, text: str, base: int, min_size: int = 12, threshold: int = 44) -> int:
        if not text:
            return base
        overflow = max(0, len(text.strip()) - threshold)
        steps = (overflow + 19) // 20
        return max(min_size, base - steps * 2)

    def _add_card(self, slide, left: float, top: float, width: float, height: float):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.theme["surface"]
        card.line.fill.background()
        return card

    def _add_progress_footer(self, slide, slide_number: int) -> None:
        if self.total_slides <= 1:
            return
        self._add_text(slide, 0.9, 6.86, 1.2, 0.18, f"{slide_number:02d}/{self.total_slides:02d}", 10, self.theme["muted"], bold=True)
        rail = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.0),
            Inches(6.92),
            Inches(9.8),
            Inches(0.10),
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = self.theme["surface"]
        rail.line.fill.background()
        progress_width = max(0.15, 9.8 * (slide_number / self.total_slides))
        fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.0),
            Inches(6.92),
            Inches(progress_width),
            Inches(0.10),
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = self.theme["accent"]
        fill.line.fill.background()

    def _is_chinese(self, text: str) -> bool:
        return any("\u4e00" <= char <= "\u9fff" for char in text)

    def _render_cover(self, slide, slide_spec: SlideSpec) -> None:
        self._add_orb(slide, 9.7, 0.5, 2.5, self.theme["accent"], 0.72)
        self._add_orb(slide, 10.8, 4.8, 1.4, self.theme["secondary"], 0.78)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.9),
            Inches(1.5),
            Inches(0.10),
            Inches(3.9),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme["accent"]
        accent_bar.line.fill.background()
        self._add_text(
            slide,
            1.25,
            1.65,
            8.1,
            2.5,
            slide_spec.title,
            self._fit_font(slide_spec.title, 38, min_size=26, threshold=26),
            self.theme["primary"],
            font_name=self.theme["display_font"],
            bold=True,
            line_spacing=1.06,
        )
        if slide_spec.thesis:
            self._add_text(slide, 1.28, 4.35, 7.5, 1.15, slide_spec.thesis, 18, self.theme["secondary"], line_spacing=1.14)

    def _render_agenda(self, slide, slide_spec: SlideSpec) -> None:
        count = max(1, len(slide_spec.bullets))
        if self._is_chinese(slide_spec.title):
            headline = f"这份演示分成 {count} 个部分。"
        else:
            headline = f"This deck covers {count} sections."
        self._add_text(slide, 0.9, 1.02, 11.0, 0.55, slide_spec.title, 24, self.theme["primary"], bold=True)
        self._add_text(slide, 0.9, 1.58, 10.2, 0.48, headline, 16, self.theme["secondary"], bold=False)

        # Support up to 6 items in 3x2 grid layout
        bullets = slide_spec.bullets[:6]
        card_width = 3.8
        card_height = 1.35
        col_gap = 0.4
        row_gap = 0.3
        start_left = 0.9
        start_top = 2.2

        for index, bullet in enumerate(bullets):
            col = index % 3
            row = index // 3
            left = start_left + col * (card_width + col_gap)
            top = start_top + row * (card_height + row_gap)

            self._add_card(slide, left, top, card_width, card_height)
            self._add_text(slide, left + 0.18, top + 0.14, 0.6, 0.18, f"{index + 1:02d}", 10, self.theme["accent"], bold=True)

            # Smaller font for more items, fit to card
            font_size = self._fit_font(bullet, 14, min_size=10, threshold=28)

            self._add_text(
                slide,
                left + 0.18,
                top + 0.36,
                card_width - 0.36,
                card_height - 0.54,
                bullet,
                font_size,
                self.theme["primary"],
                bold=True,
                line_spacing=1.20,
            )

    def _render_section(self, slide, slide_spec: SlideSpec) -> None:
        self._add_orb(slide, 9.8, 0.8, 2.0, self.theme["accent"], 0.78)
        self._add_text(
            slide,
            0.9,
            1.85,
            9.2,
            2.0,
            slide_spec.title,
            self._fit_font(slide_spec.title, 36, min_size=24, threshold=32),
            self.theme["primary"],
            font_name=self.theme["display_font"],
            bold=True,
            line_spacing=1.04,
        )
        self._add_text(slide, 0.9, 4.2, 8.5, 1.0, slide_spec.thesis, 18, self.theme["secondary"], line_spacing=1.12)

    def _render_point_like(self, slide, slide_spec: SlideSpec) -> None:
        self._add_text(slide, 0.9, 1.0, 6.3, 0.6, slide_spec.title, self._fit_font(slide_spec.title, 26, min_size=18, threshold=34), self.theme["primary"], bold=True)
        self._add_text(
            slide,
            0.9,
            1.95,
            6.2,
            2.35,
            slide_spec.thesis,
            self._fit_font(slide_spec.thesis, 24, min_size=16, threshold=40),
            self.theme["primary"],
            font_name=self.theme["display_font"],
            bold=True,
            line_spacing=1.15,
        )
        self._add_card(slide, 7.25, 1.18, 5.0, 5.25)
        # 统一列表字体大小
        bullets = slide_spec.bullets[:4]
        bullet_font_size = 12
        if bullets:
            bullet_font_size = min(12, self._fit_font(max(bullets, key=len), 12, min_size=10, threshold=40))
        for index, bullet in enumerate(bullets):
            top = 1.65 + index * 1.25
            self._add_card(slide, 7.55, top, 4.35, 1.05)
            self._add_text(
                slide,
                7.82,
                top + 0.10,
                3.85,
                0.88,
                bullet,
                bullet_font_size,
                self.theme["primary"],
                line_spacing=1.30,
            )

    def _render_compare(self, slide, slide_spec: SlideSpec) -> None:
        self._add_text(slide, 0.9, 1.0, 10.6, 0.62, slide_spec.title, self._fit_font(slide_spec.title, 26, min_size=18, threshold=36), self.theme["primary"], bold=True)
        self._add_text(slide, 0.9, 1.72, 10.4, 0.6, slide_spec.thesis, 16, self.theme["secondary"], line_spacing=1.2)
        self._add_card(slide, 0.9, 2.55, 5.55, 3.7)
        self._add_card(slide, 6.9, 2.55, 5.55, 3.7)
        left = slide_spec.bullets[:2]
        right = slide_spec.bullets[2:4] or slide_spec.bullets[:2]
        for index, item in enumerate(left):
            self._add_text(slide, 1.2, 2.95 + index * 1.24, 4.95, 0.9, f"• {item}", 15, self.theme["primary"], line_spacing=1.25)
        for index, item in enumerate(right):
            self._add_text(slide, 7.2, 2.95 + index * 1.24, 4.95, 0.9, f"• {item}", 15, self.theme["primary"], line_spacing=1.25)

    def _render_process(self, slide, slide_spec: SlideSpec) -> None:
        self._add_text(slide, 0.9, 1.0, 10.4, 0.65, slide_spec.title, self._fit_font(slide_spec.title, 25, min_size=18, threshold=34), self.theme["primary"], bold=True)
        rail = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.3),
            Inches(2.45),
            Inches(0.07),
            Inches(3.05),
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = self.theme["accent"]
        rail.line.fill.background()
        for index, item in enumerate(slide_spec.bullets[:4], start=1):
            top = 2.18 + (index - 1) * 0.98
            self._add_orb(slide, 1.05, top + 0.02, 0.34, self.theme["accent"], 0.15)
            self._add_card(slide, 1.72, top, 9.8, 0.72)
            self._add_text(slide, 1.11, top + 0.09, 0.18, 0.14, str(index), 12, self.theme["background"], bold=True, align=PP_ALIGN.CENTER)
            self._add_text(slide, 2.0, top + 0.10, 9.1, 0.54, item, 14, self.theme["primary"], line_spacing=1.20)

    def _render_quote(self, slide, slide_spec: SlideSpec) -> None:
        self._add_orb(slide, 9.9, 1.0, 1.9, self.theme["accent"], 0.18)
        self._add_text(slide, 0.78, 1.0, 0.8, 0.8, "“", 46, self.theme["accent"], font_name=self.theme["display_font"])
        self._add_text(
            slide,
            1.58,
            1.52,
            10.2,
            2.4,
            slide_spec.thesis or slide_spec.body or slide_spec.title,
            26,
            self.theme["primary"],
            font_name=self.theme["display_font"],
            bold=True,
        )
        self._add_text(slide, 1.6, 4.6, 8.8, 0.25, slide_spec.title, 16, self.theme["secondary"])

    def _render_image(self, slide, slide_spec: SlideSpec) -> None:
        self._add_text(slide, 0.9, 1.0, 10.5, 0.64, slide_spec.title, self._fit_font(slide_spec.title, 25, min_size=18, threshold=34), self.theme["primary"], bold=True)
        if slide_spec.image_path and Path(slide_spec.image_path).exists():
            slide.shapes.add_picture(str(slide_spec.image_path), Inches(0.9), Inches(1.65), width=Inches(5.8))
        else:
            self._add_card(slide, 0.9, 1.65, 5.8, 4.25)
        self._add_text(
            slide,
            7.08,
            1.82,
            5.0,
            1.4,
            slide_spec.thesis,
            self._fit_font(slide_spec.thesis, 21, min_size=14, threshold=46),
            self.theme["secondary"],
            line_spacing=1.20,
        )
        for index, bullet in enumerate(slide_spec.bullets[:3]):
            self._add_text(slide, 7.08, 3.25 + index * 0.95, 4.8, 0.85, f"• {bullet}", 14, self.theme["primary"], line_spacing=1.25)

    def _render_data(self, slide, slide_spec: SlideSpec) -> None:
        self._add_text(slide, 0.9, 1.0, 11.0, 0.62, slide_spec.title, self._fit_font(slide_spec.title, 25, min_size=18, threshold=36), self.theme["primary"], bold=True)
        hero = slide_spec.bullets[0] if slide_spec.bullets else slide_spec.thesis
        self._add_text(
            slide,
            0.9,
            1.95,
            6.0,
            1.28,
            hero,
            self._fit_font(hero, 28, min_size=17, threshold=38),
            self.theme["accent"],
            bold=True,
            font_name=self.theme["display_font"],
            line_spacing=1.06,
        )
        self._add_text(slide, 0.9, 3.45, 5.9, 0.88, slide_spec.thesis, 15, self.theme["secondary"], line_spacing=1.2)
        for index, bullet in enumerate(slide_spec.bullets[1:4], start=1):
            self._add_card(slide, 7.2, 1.92 + (index - 1) * 1.18, 4.8, 0.84)
            self._add_text(slide, 7.5, 2.10 + (index - 1) * 1.18, 4.1, 0.56, bullet, 14, self.theme["primary"], line_spacing=1.25)

    def _render_closing(self, slide, slide_spec: SlideSpec) -> None:
        self._add_orb(slide, 10.0, 0.9, 1.8, self.theme["accent"], 0.18)
        self._add_text(
            slide,
            2.15,
            2.68,
            9.0,
            1.0,
            slide_spec.title,
            40,
            self.theme["primary"],
            font_name=self.theme["display_font"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if slide_spec.thesis:
            self._add_text(slide, 2.2, 4.0, 8.9, 0.6, slide_spec.thesis, 18, self.theme["secondary"], align=PP_ALIGN.CENTER)

    def render_slide(self, slide_spec: SlideSpec, slide_number: int) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        layout = slide_spec.layout_hint

        if layout == "cover-hero":
            self._render_cover(slide, slide_spec)
        elif layout == "agenda-grid":
            self._render_agenda(slide, slide_spec)
        elif layout == "section-divider":
            self._render_section(slide, slide_spec)
        elif layout == "compare-columns":
            self._render_compare(slide, slide_spec)
        elif layout == "process-track":
            self._render_process(slide, slide_spec)
        elif layout == "quote-poster":
            self._render_quote(slide, slide_spec)
        elif layout == "image-split":
            self._render_image(slide, slide_spec)
        elif layout == "data-dashboard":
            self._render_data(slide, slide_spec)
        elif layout == "closing-clean":
            self._render_closing(slide, slide_spec)
        else:
            self._render_point_like(slide, slide_spec)

        self._add_progress_footer(slide, slide_number)

    def save(self, deck: DeckSpec, output_path: Path) -> Path:
        for slide_number, slide_spec in enumerate(deck.slides, start=1):
            self.render_slide(slide_spec, slide_number)
        ensure_parent_dir(output_path)
        self.prs.save(str(output_path))
        return output_path


def render_deck_to_pptx(deck: DeckSpec, output_path: Path) -> Path:
    renderer = PPTXRenderer(deck.style_direction, total_slides=len(deck.slides))
    return renderer.save(deck, output_path)
