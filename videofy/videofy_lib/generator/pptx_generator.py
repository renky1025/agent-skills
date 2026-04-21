"""
PPTX 生成器 - 使用 python-pptx 生成 PowerPoint
"""

import os
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

from ..models import Slide
from ..config import THEMES, PPTX_WIDTH, PPTX_HEIGHT


class PPTXGenerator:
    """PPTX 幻灯片生成器"""

    def __init__(self, theme: str = "linear"):
        self.theme = THEMES.get(theme, THEMES["linear"])
        self.prs = Presentation()
        # 设置幻灯片尺寸为 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _rgb(self, color_key: str) -> RGBColor:
        """从主题获取 RGB 颜色"""
        hex_color = self.theme.get(color_key, "#FFFFFF")
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _add_background(self, slide):
        """添加深色背景"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb("bg_dark")
        bg.line.fill.background()
        # 将背景移到最底层
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_text_box(self, slide, left, top, width, height, text: str,
                      font_size: int = 18, bold: bool = False,
                      color_key: str = "text_primary", align=PP_ALIGN.LEFT):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = self._rgb(color_key)
        p.alignment = align

        return txBox

    def _add_bullet_points(self, slide, left, top, width, height, bullets: List[str],
                           font_size: int = 16, color_key: str = "text_secondary"):
        """添加项目符号列表"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = self._rgb(color_key)
            p.space_after = Pt(12)

        return txBox

    def _add_accent_bar(self, slide, left, top, width, height):
        """添加强调色条"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        bar.fill.solid()
        # 使用 accent_rgb 或 fallback
        if "accent_rgb" in self.theme:
            r, g, b = self.theme["accent_rgb"]
            bar.fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            bar.fill.fore_color.rgb = self._rgb("accent")
        bar.line.fill.background()
        return bar

    def _render_title(self, slide_data: Slide):
        """渲染标题页"""
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 添加装饰线
        self._add_accent_bar(slide, 1.5, 2.8, 0.08, 1.5)

        # 主标题
        self._add_text_box(
            slide, 1.8, 2.5, 10, 1.5,
            content.title,
            font_size=54, bold=True,
            color_key="text_primary"
        )

        # 副标题
        if content.subtitle:
            self._add_text_box(
                slide, 1.8, 4.2, 10, 0.8,
                content.subtitle,
                font_size=24,
                color_key="text_secondary"
            )

        return slide

    def _render_agenda(self, slide_data: Slide):
        """渲染目录页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 标题
        self._add_text_box(
            slide, 1, 0.8, 11, 0.8,
            "目录",
            font_size=36, bold=True,
            color_key="text_primary"
        )

        # 装饰线
        self._add_accent_bar(slide, 1, 1.6, 1.5, 0.05)

        # 目录项
        bullets = content.bullets[:6]  # 最多显示6个
        self._add_bullet_points(
            slide, 1.5, 2.2, 10, 4.5,
            [f"{i+1}. {b}" for i, b in enumerate(bullets)],
            font_size=22, color_key="text_secondary"
        )

        return slide

    def _render_topic_cover(self, slide_data: Slide):
        """渲染主题封面页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 章节编号
        if slide_data.agenda_item:
            self._add_text_box(
                slide, 1, 2.5, 2, 1,
                f"0{slide_data.agenda_item.index}" if slide_data.agenda_item.index < 10 else str(slide_data.agenda_item.index),
                font_size=72, bold=True,
                color_key="accent"
            )

        # 主题标题
        self._add_text_box(
            slide, 1, 3.5, 11, 1.2,
            content.title,
            font_size=48, bold=True,
            color_key="text_primary"
        )

        # 摘要
        if content.subtitle:
            self._add_text_box(
                slide, 1, 4.8, 10, 1,
                content.subtitle,
                font_size=20,
                color_key="text_secondary"
            )

        return slide

    def _render_topic_overview(self, slide_data: Slide):
        """渲染主题概览页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 标签
        self._add_text_box(
            slide, 1, 0.6, 3, 0.5,
            "Overview",
            font_size=14, bold=True,
            color_key="accent"
        )

        # 标题
        self._add_text_box(
            slide, 1, 1, 11, 0.8,
            content.title,
            font_size=36, bold=True,
            color_key="text_primary"
        )

        # 要点卡片
        bullets = content.bullets[:4]
        if bullets:
            self._add_bullet_points(
                slide, 1.5, 2.2, 10, 4.5,
                bullets,
                font_size=20, color_key="text_secondary"
            )

        return slide

    def _render_topic_detail(self, slide_data: Slide):
        """渲染主题详情页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 标题
        self._add_text_box(
            slide, 1, 0.8, 11, 0.8,
            content.title,
            font_size=36, bold=True,
            color_key="text_primary"
        )

        # 装饰线
        self._add_accent_bar(slide, 1, 1.6, 1, 0.05)

        # 要点列表
        bullets = content.bullets[:6]
        if bullets:
            self._add_bullet_points(
                slide, 1.5, 2.2, 10, 4.5,
                bullets,
                font_size=18, color_key="text_secondary"
            )

        return slide

    def _render_split_left(self, slide_data: Slide):
        """渲染左图右文布局"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 左侧图片 (40%)
        if content.image_path and os.path.exists(content.image_path):
            try:
                slide.shapes.add_picture(
                    content.image_path,
                    Inches(0.5), Inches(1.5),
                    width=Inches(5)
                )
            except Exception:
                pass  # 图片加载失败时跳过

        # 右侧文字 (60%)
        self._add_text_box(
            slide, 6, 1, 7, 0.8,
            content.title,
            font_size=32, bold=True,
            color_key="text_primary"
        )

        # 图片说明/正文
        if content.body_text:
            self._add_text_box(
                slide, 6, 2.2, 6.5, 4,
                "\n".join(content.body_text),
                font_size=16,
                color_key="text_secondary"
            )

        return slide

    def _render_split_right(self, slide_data: Slide):
        """渲染右图左文布局"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 左侧文字 (60%)
        self._add_text_box(
            slide, 0.5, 1, 7, 0.8,
            content.title,
            font_size=32, bold=True,
            color_key="text_primary"
        )

        # 图片说明/正文
        if content.body_text:
            self._add_text_box(
                slide, 0.5, 2.2, 6.5, 4,
                "\n".join(content.body_text),
                font_size=16,
                color_key="text_secondary"
            )

        # 右侧图片 (40%)
        if content.image_path and os.path.exists(content.image_path):
            try:
                slide.shapes.add_picture(
                    content.image_path,
                    Inches(7), Inches(1.5),
                    width=Inches(5.5)
                )
            except Exception:
                pass

        return slide

    def _render_code(self, slide_data: Slide):
        """渲染代码页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 标题
        self._add_text_box(
            slide, 0.5, 0.5, 12, 0.6,
            content.title,
            font_size=24, bold=True,
            color_key="text_primary"
        )

        # 代码背景
        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.3), Inches(5.8)
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = self._rgb("bg_code")
        code_bg.line.fill.background()

        # 语言标签
        if content.code_lang:
            self._add_text_box(
                slide, 0.7, 1.3, 2, 0.4,
                content.code_lang,
                font_size=12,
                color_key="text_muted"
            )

        # 代码内容
        if content.code:
            code_text = content.code[:800]  # 限制长度
            self._add_text_box(
                slide, 0.7, 1.7, 11.8, 5,
                code_text,
                font_size=11,
                color_key="text_secondary"
            )

        return slide

    def _render_end(self, slide_data: Slide):
        """渲染结束页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background(slide)

        content = slide_data.content

        # 居中文本
        self._add_text_box(
            slide, 3, 3, 7.333, 1.5,
            content.title,
            font_size=54, bold=True,
            color_key="text_primary",
            align=PP_ALIGN.CENTER
        )

        return slide

    def generate(self, slides: List[Slide]) -> Presentation:
        """生成完整 PPTX"""
        for slide_data in slides:
            layout_type = slide_data.layout_type

            if layout_type == "title":
                self._render_title(slide_data)
            elif layout_type == "agenda":
                self._render_agenda(slide_data)
            elif layout_type == "topic_cover":
                self._render_topic_cover(slide_data)
            elif layout_type == "topic_overview":
                self._render_topic_overview(slide_data)
            elif layout_type == "topic_detail":
                self._render_topic_detail(slide_data)
            elif layout_type == "split_left":
                self._render_split_left(slide_data)
            elif layout_type == "split_right":
                self._render_split_right(slide_data)
            elif layout_type == "code":
                self._render_code(slide_data)
            elif layout_type == "end":
                self._render_end(slide_data)
            else:
                # 默认使用 topic_detail
                self._render_topic_detail(slide_data)

        return self.prs

    def save(self, slides: List[Slide], output_path: str):
        """保存 PPTX 到文件"""
        self.generate(slides)
        self.prs.save(output_path)
        return output_path
